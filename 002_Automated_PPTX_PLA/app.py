'''
App made by:
    Eduardo Reyes Alvarez, Ph.D.
Contact:
    eduardo_reyes09@hotmail.com

App version: 
    V03 (Nov 02, 2023): Reordered parts of the code to make it more organized. Improved app layout
                        and transferred info text from the Colab notebook version to here.

'''
###################################################################################################

# Import required libraries

import os
import glob
from zipfile import ZipFile
import urllib.request
import pandas as pd
import time
import streamlit as st
from streamlit_option_menu import option_menu

# Python-pptx specific modules

# To make the presentation
from pptx import Presentation
# To specify sizes of images and text (Other options like inches are available)
from pptx.util import Cm, Pt            
# To specify text alignment (the method is used on the text frame, not on the box or the actual text)
from pptx.enum.text import PP_ALIGN     
# To work with line and fill colours
from pptx.enum.dml import MSO_THEME_COLOR

###################################################################################################

# App main layout and page loading

st.set_page_config(
    page_title="Tool 01 - App by Eduardo",
    page_icon=":newspaper:",
    layout="wide")

st.title("Pptx generator for PLA results")
st.write("---")

# Make a menu of pages on the siderbar, since the app is simple but requires lots of specific details
with st.sidebar:
    selected_page = option_menu("Main Menu", ["Generate pptx", "How to use the app", "Info on pptx design"], 
        icons=["filetype-pptx", "patch-question-fill", "columns-gap"], menu_icon="cast", default_index=0)

# Check the selected app page and call the corresponding function to display its content
def change_pages():
    if selected_page == "Generate pptx":
        load_first_page()
    elif selected_page == "How to use the app":
        load_second_page()
    else:
        load_third_page() 

    # Get the working directory for all functions to write files to it, and download the blank pptx template
    st.session_state["current_directory"] = os.path.dirname(os.path.realpath(__file__))
    st.session_state["template_pptx"] = os.path.join(st.session_state["current_directory"], "Template.pptx")
    if not os.path.exists(st.session_state["template_pptx"]):
        template_dir = "https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Template.pptx"
        urllib.request.urlretrieve(template_dir, st.session_state["template_pptx"])

###################################################################################################

# Function to process the input files

@st.cache_data
def process_files():

    # Extract all the folder structure and files of the zip file provided by the user
    with ZipFile(st.session_state["data_zipfile"], 'r') as zip:
        zip.extractall()

    # The unzipping takes several seconds, so we wait for the folder to appear
    while not "Data" in os.listdir():
        time.sleep(1)

    # Walk through the Data folder to find all csv files
    data_folder = os.path.join(st.session_state["current_directory"], "Data")
    extension_wanted = "*.csv"
    all_csv_files = [file
                    for path, subdir, files in os.walk(data_folder)
                    for file in glob.glob(os.path.join(path, extension_wanted))]

    # Make list with names and paths of the experimental conditions (according to the number of csv files found above)
    exp_conditions_info = []
    for csv_file in all_csv_files:
        csv_root_folder = csv_file.split("Data/")[1].split("/Quantification")[0]
        exp_conditions_info.append([csv_root_folder, csv_file])

    # We will store all the information of the images to insert to the ppt here
    all_info_for_slides = []

    # Iterate for each folder contained in the Data folder uploaded
    for condition_info in exp_conditions_info:
        
        # Read the csv file for the current experimental condition
        results_table = pd.read_csv(condition_info[1])

        # Make some edits for easier manipulation and sorting of the subtitles, folders and ROI names
        results_table["Image used"] = results_table["Image used"].str.replace("MAX_","")
        results_table["Image used"] = results_table["Image used"].str.replace(".tif","")
        results_table["Cell quantified"] = results_table["Cell quantified"].str.replace("_1.roi","")
        results_table["Cell quantified"] = results_table["Cell quantified"].apply(int)
        results_table["Particle count threshold"] = results_table["Particle count threshold"].apply(int)
        results_table["Particle count maxima"] = results_table["Particle count maxima"].apply(int)
        results_table = results_table.sort_values(by=["Image used", "Cell quantified"],ignore_index=True)

        # Iterate through the rows of the csv file (ROIs/cells quantified)
        for index, row in results_table.iterrows():
            
            # Retrieve all the relevant information for that row
            ROI_title = condition_info[0]
            ROI_subtitle = row["Image used"]
            ROI_name = str(row["Cell quantified"])
            ROI_Fluorescence_image = condition_info[1].replace("Quantification/Results.csv", "Cropped cells/Fluorescence/"+ROI_subtitle+"/"+ROI_name+"_2.jpg")
            ROI_Tcount_image = ROI_Fluorescence_image.replace("Fluorescence", "T_Particles").replace("_2.jpg", "_1.jpg")
            ROI_Tcount = row["Particle count threshold"]
            ROI_FMcount_image = ROI_Fluorescence_image.replace("Fluorescence", "FM_Particles").replace("_2.jpg", "_1.jpg")
            ROI_FMcount = row["Particle count maxima"]

            all_info_for_slides.append([ROI_title, ROI_subtitle, ROI_Fluorescence_image, ROI_name, ROI_Tcount_image, ROI_Tcount, ROI_FMcount_image, ROI_FMcount])

    # Prepare empty variables to store the grouped images for each slide
    all_slides_content = []
    temp_slide = []

    # Retrieve the title and subtitle of the very first image to feed the loop at i=0
    current_title = all_info_for_slides[0][0]
    current_subtitle = all_info_for_slides[0][1]

    # Iterate through all the information retrieve for each cell/ROI
    for info in all_info_for_slides:
        
        # Get the current title and subtitle to compare to the reference
        new_title = info[0]
        new_subtitle = info[1]

        # Check the 3 conditions described in text above
        if len(temp_slide)==20 or new_title!=current_title or new_subtitle!=current_subtitle:
            
            # If anything triggers the change of slide, dump the current, empty it and set titles as refs
            all_slides_content.append(temp_slide)
            temp_slide = []
            current_title = new_title
            current_subtitle = new_subtitle
        
        # Always attach the cell/ROI, could be at the same group/slide with others, or to a empty temp slide
        temp_slide.append(info)

    return all_slides_content

###################################################################################################

# Function to generate the presentations and pass the slide maker the content it should insert 

def generate_pptxs(all_slides_content):

    # This will make two presentations, one for the Thresholding and one for the Find Maxima methods

    # Open the presentation uploaded by the user
    presentation_T = Presentation(st.session_state["template_pptx"])

    # Iterate through the image info grouped by slide
    for slide_content in all_slides_content:

        # Prepare the parameters we need to pass to the function that makes the slides
        current_slide_title = slide_content[0][0]
        current_slide_subtitle = slide_content[0][1]
        image_count_for_slide = len(slide_content)
        F_images_for_slide = [image[2] for image in slide_content]
        F_images_labels = [image[3] for image in slide_content]
        P_images_for_slide = [image[4] for image in slide_content]
        P_images_labels = [image[5] for image in slide_content]

        # Feed the function that makes the slide and inserts the corresponding images
        presentation_T = slide_maker(presentation_T, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                                     F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels)

    # Finally, save this summary presentation after all slides have been created
    presentation_T.save(os.path.join(st.session_state["current_directory"], "Summary_results_T.pptx"))

    # Open the presentation uploaded by the user
    presentation_FM = Presentation(st.session_state["template_pptx"])

    # Iterate through the image info grouped by slide
    for slide_content in all_slides_content:

        # Prepare the parameters we need to pass to the function that makes the slides
        current_slide_title = slide_content[0][0]
        current_slide_subtitle = slide_content[0][1]
        image_count_for_slide = len(slide_content)
        F_images_for_slide = [image[2] for image in slide_content]
        F_images_labels = [image[3] for image in slide_content]
        P_images_for_slide = [image[6] for image in slide_content]
        P_images_labels = [image[7] for image in slide_content]
        
        # Feed the function
        presentation_FM = slide_maker(presentation_FM, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                                      F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels)

    # Finally, save this summary presentation after all slides have been created
    presentation_FM.save(os.path.join(st.session_state["current_directory"], "Summary_results_FM.pptx"))

###################################################################################################

# Function to add slides to a pptx and inserts images+text 

def slide_maker(presentation_input, current_slide_title, current_slide_subtitle, image_count_for_slide, 
                F_images_for_slide, P_images_for_slide, F_images_labels, P_images_labels):

    # All coordinates are stated always in the same order: From left first, from top second.

    # Title text box dimensions and coordinates (centimeters) 
    title_width = 17
    title_height = 1.5
    title_left_coordinate = 0
    title_top_coordinate = 0
    
    # Subtitle text box dimensions and coordinates (centimeters)
    subtitle_width = 17
    subtitle_height = 1.5
    subtitle_left_coordinate = 17
    subtitle_top_coordinate = 0
    
    # Size and coordinates for the 20 pairs of images (centimeters)
    image_width = 3.25
    image_height = 3
    image_coordinates = [
    (0.25, 2.1, 3.5, 2.1),   (7, 2.1, 10.25, 2.1),   (13.75 , 2.1, 17, 2.1),  (20.5, 2.1, 23.75, 2.1),    (27.25, 2.1, 30.5, 2.1),
    (0.25, 6.4, 3.5, 6.4),   (7, 6.4, 10.25, 6.4),   (13.75, 6.4, 17, 6.4),   (20.5, 6.4, 23.75, 6.4),    (27.25, 6.4, 30.5, 6.4),
    (0.25, 10.7, 3.5, 10.7), (7, 10.7, 10.25, 10.7), (13.75, 10.7, 17, 10.7), (20.5, 10.7, 23.75, 10.7),  (27.25, 10.7, 30.5, 10.7),
    (0.25, 15, 3.5, 15),     (7, 15, 10.25, 15),     (13.75, 15, 17, 15),     (20.5, 15, 23.75, 15),      (27.25, 15, 30.5, 15)
    ]
    
    # Size and coordinates for the 20 pairs of text labels (centimeters) (+3cm top coordinate of images)
    image_labels_width = 3.25
    image_labels_height = 1
    image_labels_coordinates = [
    (0.25, 5.1, 3.5, 5.1),   (7, 5.1, 10.25, 5.1),   (13.75 , 5.1, 17, 5.1),  (20.5, 5.1, 23.75, 5.1),    (27.25, 5.1, 30.5, 5.1),
    (0.25, 9.4, 3.5, 9.4),   (7, 9.4, 10.25, 9.4),   (13.75, 9.4, 17, 9.4),   (20.5, 9.4, 23.75, 9.4),    (27.25, 9.4, 30.5, 9.4),
    (0.25, 13.7, 3.5, 13.7), (7, 13.7, 10.25, 13.7), (13.75, 13.7, 17, 13.7), (20.5, 13.7, 23.75, 13.7),  (27.25, 13.7, 30.5, 13.7),
    (0.25, 18, 3.5, 18),     (7, 18, 10.25, 18),     (13.75, 18, 17, 18),     (20.5, 18, 23.75, 18),      (27.25, 18, 30.5, 18)
    ]
    
    # Create a new slide (layout Blank)
    blank_slide_layout = presentation_input.slide_layouts[6]
    slide = presentation_input.slides.add_slide(blank_slide_layout)
    
    # Make the title for this experimental condition
    left = Cm(title_left_coordinate)
    top = Cm(title_top_coordinate)
    width = Cm(title_width)
    height = Cm(title_height)
    title_textbox = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_textbox.text_frame
    title_text = title_frame.paragraphs[0]
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_text.text = current_slide_title
    title_text.font.bold = True
    title_text.font.size = Pt(32)
    title_text.font.name = "Times New Roman"
    
    # Make the subtitle for the image where the ROI was cropped from
    left = Cm(subtitle_left_coordinate)
    top = Cm(subtitle_top_coordinate)
    width = Cm(subtitle_width)
    height = Cm(subtitle_height)
    subtitle_textbox = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_textbox.text_frame
    subtitle_text = subtitle_frame.paragraphs[0]
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_text.text = current_slide_subtitle
    subtitle_text.font.size = Pt(32)
    subtitle_text.font.name = "Times New Roman"
    
    # Based on the number of images for the current slide, retrieve the neccesary images and coordinates
    for i in range(image_count_for_slide):
        
        # Find the images to insert
        fluorescence_image = F_images_for_slide[i]
        particle_image = P_images_for_slide[i]
        
        # Insert the cropped cell from the Fluorescence folder first
        left = Cm(image_coordinates[i][0])
        top = Cm(image_coordinates[i][1])
        width = Cm(image_width)
        height = Cm(image_height)
        inserting_image = slide.shapes.add_picture(fluorescence_image, left, top, width, height)
        
        # Insert the text label corresponding to the image just inserted above
        left = Cm(image_labels_coordinates[i][0])
        top = Cm(image_labels_coordinates[i][1])
        width = Cm(image_labels_width)
        height = Cm(image_labels_height)
        inserting_image_textbox = slide.shapes.add_textbox(left, top, width, height)
        inserting_image_frame = inserting_image_textbox.text_frame
        inserting_image_text = inserting_image_frame.paragraphs[0]
        inserting_image_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        inserting_image_text.text = F_images_labels[i]
        inserting_image_text.font.size = Pt(20)
        inserting_image_text.font.name = "Times New Roman"
        
        # Insert the cropped cell from the Particles folder second (FM or T particles)
        left = Cm(image_coordinates[i][2])
        top = Cm(image_coordinates[i][3])
        width = Cm(image_width)
        height = Cm(image_height)
        inserting_image2 = slide.shapes.add_picture(particle_image, left, top, width, height)
        inserting_image2.line.fill.solid()
        inserting_image2.line.width = Pt(0.5)
        inserting_image2.line.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
        # Insert the text label corresponding to the particle counts just inserted above
        left = Cm(image_labels_coordinates[i][2])
        top = Cm(image_labels_coordinates[i][3])
        width = Cm(image_labels_width)
        height = Cm(image_labels_height)
        inserting_image2_textbox = slide.shapes.add_textbox(left, top, width, height)
        inserting_image2_frame = inserting_image2_textbox.text_frame
        inserting_image2_text = inserting_image2_frame.paragraphs[0]
        inserting_image2_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        inserting_image2_text.text = "P="+str(P_images_labels[i])
        inserting_image2_text.font.size = Pt(20)
        inserting_image2_text.font.name = "Times New Roman"
    
    return presentation_input

###################################################################################################

# Function to load the first app page which takes the input Data.zip file and producess the outputs

def load_first_page():
    
    # Show the user a widget to upload the compressed file
    uploaded_file = st.file_uploader("Upload compressed file", type=["zip"], accept_multiple_files=False)
    st.write("---")

    #Create columns for better layout of the buttons
    col_1_row_1, col_2_row_1, col_3_row_1 = st.columns([1, 1, 1], gap="medium")

    # Display a button so the user decides when to start (in case uploaded the incorrect file)
    with col_1_row_1:
        start_button = st.button(label="Generate pptx", type="primary")
    
    # Proceed only when the button to start is pressed and a compressed file has been uploaded
    if start_button and uploaded_file:
        
        with open(uploaded_file.name, "wb") as f:
            f.write(uploaded_file.getvalue())
        st.session_state["data_zipfile"] = os.path.join(os.path.dirname(os.path.abspath(uploaded_file.name)), "Data.zip")

        # Process the files to extract the information needed to import to the slide generator
        all_slides_content = process_files()
        generate_pptxs(all_slides_content)

        # Display download buttons for each of the two files produced
        threholding_file_path = os.path.join(st.session_state["current_directory"], "Summary_results_T.pptx")
        find_maxima_file_path = os.path.join(st.session_state["current_directory"], "Summary_results_FM.pptx")
        with col_2_row_1:
            st.download_button(label="Download Threholding file", key="download_threholding", 
                               data=open(threholding_file_path, "rb").read(), file_name="Summary_results_T.pptx")
        with col_3_row_1:
            st.download_button(label="Download Find Maxima file", key="download_find_maxima", 
                               data=open(find_maxima_file_path, "rb").read(), file_name="Summary_results_FM.pptx")

    return

###################################################################################################

# Function to load the second app page with basic instructions on how to use the app

def load_second_page():
    
    st.markdown('''
                ## Function of the notebook

                This notebook takes any number of Proximity Ligation Assay (PLA) images that were cropped 
                and quantified, and prepares a summary Power Point presentation with the results of the 
                experiment. In the presentation we will find:

                1. Slides for each experimental condition / group (separated), and also separated based 
                on the original image they were croped from (if big images were taken with multiple cells, 
                the same condition will have multiple cells/ROIs originated from the same original image).

                2. Each slide shows pairs of fluorescence images of the cells selected with their ROI name 
                and next to it there is the corresponding image of the quantification of puncta for that 
                ROI and the label corresponds to the number of puncta given by the method used (Find Maxima 
                or Threshold).

                3. The key is that we can quickly preview in one single file all the cell morphologies we 
                picked to analyze, and verify the numbers match what we would expect just looking at the 
                puncta in the cell by eye. If the cells look too bad in the fluorescence images, we can 
                find the number that corresponds to that cell in the results file and discard it. Also, we 
                can evaluate the work of the quantification methods since all the black noise in those 
                quantification images correspond to ignored pixels (didn't meet the criteria) and only the 
                coloured ones were quantified. Controls should have barely any black or coloured particles 
                in the images (and thus low number P=), but the other conditions were we look for 
                interactions may have or not more black noise, coloured puncta, and numbers that match with 
                what our eyes can see.

                ''')
    
    st.write("---")
    
    st.markdown('''
                ## Requirements of the notebook

                Since this notebook was created specifically to generate a results report for experiments of 
                the author's PLA experiments, the requirements are just to upload a new pptx file with a 
                blank slide (just empty), making sure to click in Power Point: Design-> Slide Size -> 
                Widescreen (16:9), and the compressed zip file called "Data", which should contain all the 
                following:

                1. A folder called "Data" (1st level).

                2. Any number of subfolders (2nd level) corresponding to each experimental condition that you 
                want to include in the same summary presentation.

                3. Each 2nd level subfolder can have a unique name, but their content should have the same 
                structure: 2 folders (3rd level), one called "Cropped cells", and the other "Quantification". 
                Both of these were produced from the processed images by the author's script V03: 
                https://github.com/EdRey05/Resources_for_Mulligan_Lab/blob/b80eaf75d35665aeb4b7e60ed85685f342d9f125/Tools%20for%20PLA%20quantification/PLA_quantification.py

                4. The Cropped cells folder contains 3 folders (4th level) called "Fluorescence", "FM_Particles", 
                and "T_Particles". They all have subfolders with the names Row_01_05... to Row_45_51 on the 5th 
                level. The Fluorescence folder has ROIs (6th level) of cells in the form of "Number_2.jpg" and 
                the other 2 folders "Number_1.jpg", which refer to the same cell (one ROI used to make the 
                presentation, other to quantify).

                5. The Quantification folder (4th level) has only one csv file with the results of the 
                quantification.
    
                ''')
    
    st.write("---")

    st.markdown('''
                ## Outputs of the script

                The current version of this notebook produces 2 pptx files, one using the fluorescence images + 
                the find maxima images, and the second one using the same fluorescence images but with the 
                thresholded images. Ideally, we would want to examine both to check which works better, since 
                under ideal conditions both give very similar results, but raw image quality and processing 
                can create some differences within a condition or experiment so using the same threshold method 
                will cause a lot of artifacts in some images (find maxima uses a different principle, more 
                insensitive to these variability).
    
                ''')
    return

###################################################################################################

# Function to load the third app page which describes the design and layout of the slides

def load_third_page():
    

    st.markdown('''
            <div style='background-color: lightblue; padding: 10px; border-radius: 5px;'>

                Making this setup of automated generation of Power Point presentations is possible 
                using the ***python-pptx*** library, for which parameters need to be defined, 
                coordinates to place our objects of interest and their sizes need to be calculated. 
                For more information on this library, see:  https://pypi.org/project/python-pptx/

            </div>
                ''', unsafe_allow_html=True)


    st.markdown('''
                ## Parameters defined to insert the images on the slides

                The specific coordinates of all the elements desired in the slides were previously tested 
                on a pptx by manually arranging the images with the intended number of rows, columns, 
                images per slide, labels, text boxes and a size adequate to gain enough insight of the 
                results while keeping the number of slides as low as possible. Once all the intended 
                content for a single slide was set into approximate position, all the sizes were 
                calculated and set to precise numbers in order to make it reproducible and iterable. 

                ''')
    
    # Diplay the images of the slide coordinates
    col_1_row_1, col_2_row_1 = st.columns([1, 1], gap="small")
    with col_1_row_1:
        st.image(image="https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Automated_PPTX_goal.jpg",
                 caption="Desired slide layout", use_column_width=True,)
    with col_2_row_1:
        st.image(image="https://github.com/EdRey05/Streamlit_projects/raw/main/002_Automated_PPTX_PLA/Automated_PPTX_coordinates.jpg",
                 caption="Coordinates for each desred object", use_column_width=True)
            
    st.markdown('''
                All the parameters illustrated above are the following:

                1. Slides of 16:9 ratio (34cm width, 19cm height, all measurement units set to cm).

                2. Two titles at the very top, their text boxes of 17cm width and 1.5cm height each, 
                side by side with no separation from the top corners. Title bold font, subtitle normal 
                font, and both Times New Roman size 32 points).

                3. Below the title+subtitle, there is a 0.6cm space to the first row of images, so the 
                first row starts at X, 2.1cm (X = distance from the left edge).

                4. Each image will be resized to 3.25cm width by 3cm height. The pairs consist of a 
                fluorescence image on the left side and a particle analysis on the right side, with no 
                space separating them.

                5. There are 0.25cm separating each pair (5 pairs per row, 4 pairs per column, total 
                20 pairs/cells = 40 images per slide), and also 0.25cm separating the first and last 
                pair of a row from the edges of the slide.

                6. There is a text box right under each image (the borders touch). For the fluorescence 
                image the text refers to the name of the ROI analyzed, whereas the particle image shows 
                the count of particles for that particular cell so we can quickly inspect whether all 
                the coloured particles coincide with the actual puncta we see in the fluorescence image 
                and the final count (or if there was an error due to the quantification process). This 
                text boxes are 3.25cm width by 1cm height, contain Times New Roman text size 20 points.

                7. The second row is separated 0.3cm from the bottom edge of the labels mentioned in 6. 
                Same for the third and fourth row (which almost ends at the very edge of the slide).

                **Note:** The image and label coordinates are declared in a list of tuples (coordinates 
                can't change for this notebook), so we can iterate through the images grouped per slide 
                (1-20) and only use the required coordinates (we don't need to iterate 20 times all the 
                time if the slides only have few images, it is better to iterate per image in the slide 
                and use the index to find its corresponding coordinates). The filling order will be left 
                to right, top-down. Since the fluorescence and particle images both come from the same 
                cell, their names will be "1_2.jpg" and "1_1.jpg", respectively, so the tuple of 
                coordinates consists of 4 elements: the first is the distance from the left edge of the 
                slide for the fluorescence image, the second is the distance from the top edge of the 
                slide for the fluorescence image, the third is the distance from the left edge of the 
                slide for the particle image, and the fourth is the distance from the top edge of the 
                slide for the particle image. This way, we can use the same ROI name to get both images 
                by replacing part of the directory and the "_2.jpg" for "_1.jpg", however, the index in 
                the for loop will be the same, since we have the coordinates of both images in the same 
                element of the list (tuple of 4 coordinates).

                ''')
    
    st.write("---")

    st.markdown('''
                ## Potato

                * Unzip the "Data.zip" file to make a "Data" folder, over which we can iterate/walk 
                through.
                
                * Once the folder is ready, look for the csv results file of each experimental 
                condition (1 per folder/condition), in which we have a collection of most of the 
                information we need.

                * We know that the number of csv's found is the same number of experimental conditions 
                we need to iterate through. We get their paths, names, and we already know the folder 
                structure where the images are located: **content->Data->ExpCondition->Cropped cells->** , 
                which contains 3 subfolders: **Fluorescence, FM_Particles, and T_Particles**

                * Once we know the number of experimental condition folders, we will iterate through any 
                number of them to extract the information of all the images analyzed.

                * Since the csv files already contain almost all the information we need (subtitle, ROI 
                names, T particle could and FM particle count), we will iterate through the rows of the 
                csv file instead of walking through the directory of fluorescence images (we could also 
                extract the info from the path but we would need multiple steps to split different 
                sections of the path). This strategy also allows us to easily convert the ROI names into 
                integers so we can sort them properly (natural sorting, the ouput of the quantification 
                script is not in the correct order).

                Now we have a huge list containing all the info for all the images uploaded by the user. 
                The next step is to group these images, following a few rules:

                1. The main idea is to pass 20 images to the function that makes the slides.

                2. We start with the first experimental condition (title of slide), first subfolder of 
                the original image (like Row_01_05, which is the subtitle), and then we see how many 
                cells/ROIs were quantified there.

                3. If we have exactly 20 (unlikely), we just make one slide. If we have less than 20, we 
                also make one slide and leave empty spots. If we have more than 20 images, we take the 
                first 20, make one slide, take the next (up-to) 20 using the same title and subtitle, 
                make a new slide and so on, until we don't have more images in that subfolder.

                4. To accomplish this strategy, we need to iterate through the list generated above 
                (all_info_for_slides), in which all the information for all the images of all the 
                conditions is side by side. Because of that, we can iterate through the same level of 
                the list, take the info of the current item/element, and add its content to a new 
                temporary variable (list) with 20 spots available. Before adding the info of the current 
                item/element, we check whether the 20 spots have been filled, check whether the title 
                has changed, and check whether the subtitle has changed. Any of those 3 cases triggers 
                the jump to a new slide so we have to pass the grouped image info from the temporary 
                variable to a final variable, clear the temporary variable content and then add the 
                current image info to the new slide group. Finally, we repeat this process over and over 
                until we have checked all the items/elements, and the final variable, product of this 
                loop, will contain all the image information separated/grouped by images that will go 
                together in a single power point slide.

                5. With the final variable we will be able to iterate through each element=slide, call 
                the function that makes the slides, and pass the current element information to insert 
                the images for that slide.

                ''')
    return

###################################################################################################
change_pages()
###################################################################################################